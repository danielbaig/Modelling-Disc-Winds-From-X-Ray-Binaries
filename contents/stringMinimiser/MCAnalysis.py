import numpy as np
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches
import re
from PIL import Image

from ..pairwise import *
from . import organiseData

from pathlib import Path
pathtohere = Path.cwd()



def plotEnergyHistory(steps:np.ndarray, energyHistory_mean:np.ndarray, energyHistory_std):
    """
    Display how the energy changes with the step number.

    Inputs:
        - steps:np.ndarray: Monte Carlo step.
        - energyHistory:np.ndarray: Mean of the energy at each MC step.
        - energyHistory_std:np.ndarray: Standard deviation of the energy at each MC step.

    """

    fig = plt.figure(figsize=(8,6), dpi=200)
    ax = fig.add_subplot()

    ax.plot(steps,energyHistory_mean,c='r')
    #ax.plot(steps,energyHistory_mean+energyHistory_std,c='r',alpha=0.5)
    #ax.plot(steps,energyHistory_mean-energyHistory_std,c='r',alpha=0.5)

    # Creates appropiate labels.
    ax.set_xlabel('step')
    ax.set_ylabel('energy')

    ax.grid()

    plt.savefig(pathtohere / 'contents/stringMinimiser/plots/energyHistory.png', bbox_inches='tight')


## These will alll be integrated with pairwise.py at some point (or at leasts should).
def displayVariableParameterSpace(ax, samplePositions_x:np.ndarray,potentialPositions_x:np.ndarray,
                                        samplePositions_y:np.ndarray,potentialPositions_y:np.ndarray,
                                        actualPosition_x:np.ndarray, actualPosition_y:np.ndarray,
                                        compareSampleIndex:int):
    """
    Display a single subplot with a single point.

    Inputs:
        - ax: Instance of the subplot.
        - samplePositions_x:np.ndarray: x coordinates of the sample positions.
        - potentialPositions_x:np.ndarray: x coordinates of the potential positions.
        - samplePositions_y:np.ndarray: y coordinates of the sample positions.
        - potentialPositions_y:np.ndarray: y coordinates of the potential positions.
        - actualPosition_x:np.ndarray: x coordinate of the actual position.
        - actualPosition_y:np.ndarray: y coordinate of the actual position.
        - compareSampleIndex:int: Sample being compared to the actual position.

    Outputs:
        - ax: Modified instance of the subplot.
    """


    ax.plot(samplePositions_x,samplePositions_y, c='g',marker='.',zorder=1)
    ax.scatter(samplePositions_x[compareSampleIndex], samplePositions_y[compareSampleIndex],
                c='c',marker='.',zorder=2)
    ax.scatter(samplePositions_x[0], samplePositions_y[0], c='r', marker='x',zorder=2)
    ax.scatter(samplePositions_x[-1], samplePositions_y[-1],c='b',marker='x',zorder=2)

    ax.scatter(potentialPositions_x, potentialPositions_y,c='orange', marker='.')
    ax.scatter(actualPosition_x,actualPosition_y,c='m', marker='.')


    return ax

def _labelAxes(ax, i:int, j:int, var1:str, var2:str, numVariables:int):
    """
    Creates labels and formats the axes of a subplot.

    Inputs:
        - ax: Instance of the subplot.
        - i:int: Row number.
        - j:int: Column number.
        - var1:str: x-variable.
        - var2:str: y-variable.
        - numVariables:int: Number of variables.

    Outputs:
        - ax: Modified instance of the subplot.
    """

    # Remove tick labels on non-edge figures.
    if i!=0:
        ax.set_yticks([])
    if j!=numVariables-1:
        ax.set_xticks([])

    if i==0:
        ax.set_ylabel(var2)
    if i==j-1:
        ax.set_title(var1)

    # Put axes labels in scientific notation.
    try:
        ax.ticklabel_format(style='sci', axis='both', scilimits=(0,0))
    except:
        pass

    return ax

        

def displayPairwise(variables:tuple, samplePositions:np.ndarray,
                    potentialPositions:np.ndarray, actualPositions:np.ndarray, compareSampleIndex:int):
    """
    Display a pairwise correlation plot of the variable parameter space (for one point at the moment).

    Inputs:
        - variables:tuple: Variables used.
        - samplePositions:np.ndarray: Sample positions for the one vps point.
        - potentialPositions:np.ndarray: vps potential points.
        - actualPositions:np.ndarray: Actual position of the point.
        - compareSampleIndex:int: Sample being tested against the actual position.

    """


    numVariables = len(variables)

    fig = plt.figure(figsize=(16,16), dpi=300, tight_layout=True)

    for i,var1 in enumerate(variables):
        for j,var2 in enumerate(variables):
            if i>=j:
                continue

            ax = fig.add_subplot(numVariables,numVariables,i+j*numVariables+1)

            # Populate subplot.
            ax = displayVariableParameterSpace(ax,samplePositions[:,i],potentialPositions[:,i],
                                                samplePositions[:,j], potentialPositions[:,j],
                                                actualPositions[i], actualPositions[j], compareSampleIndex)

            # Labelling
            ax = _labelAxes(ax, i, j, var1, var2, numVariables)



    plt.savefig(pathtohere / 'contents/stringMinimiser/plots/variableParameterSpace.png', bbox_inches='tight')

##################################


def makePowerPointTitlePage(prs):
    """
    Make the title page for the powerpoint.

    Inputs:
        - prs: Instance of the presentation.

    Outputs:
        - prs: Modified instance of the presentation.
    """

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = 'Monte Carlo parameter space string minimisation prediction.'
    subtitle.text = 'Generated with python-pptx.'

    return prs


def designPowerPoint(prs, numSamples:int):
    """
    Put the images into the powerpoint.

    Inputs:
        - prs: Instance of the presentation.
        - numSamples:int: Number of samples between the known points.

    Outputs:
        - prs: Modified instance of the presentation.
    """

    blank_layout = prs.slide_layouts[6]
    x = 0
    y = 0
    height = Inches(7.5)


    for i in range(numSamples):
        slide = prs.slides.add_slide(blank_layout)
        picture = slide.shapes.add_picture(f'contents/stringMinimiser/data/alongPath/{i}.png', x, y,height=height)

    return prs


def createFrames(samplePositions:np.ndarray, variables:tuple):
    """
    Creates frames for the GIF/PowerPoint.

    Inputs:
        - samplePositions:np.ndarray: Transposed sample positions in form (variables, sample, point).
        - variables:tuple: Variables used.
    """

    # Create pairwise images.
    for i in range(samplePositions.shape[1]):
        pairwise = Pairwise(samplePositions[:,i], variables,
                            f'contents/stringMinimiser/data/alongPath/{i}.png', colours='log(t_e)',
                            directoryIsPath=True, isUpperOn=False)
        pairwise.createPairwiseFigure()




def createPowerPoint(numSamples:int):
    """
    Create a powerpoint showing how the points move through the vps.

    Inputs:
        - numSamples:int: Number of samples.
    """

    
    prs = Presentation()
    prs = makePowerPointTitlePage(prs)
    prs = designPowerPoint(prs, numSamples)
    prs.save(pathtohere / 'contents/stringMinimiser/plots/pointTrajectories.pptx')
    



def createGIF(numSamples:int):
    """
    Create a GIF of the images for the point trajectories.
    
    Inputs:
        - numSamples:int: Number of samples.
    """

    frames = [Image.open(pathtohere / f'contents/stringMinimiser/data/alongPath/{i}.png') for i in range(numSamples)]
    frames_init = frames[0]
    frames_init.save(pathtohere / 'contents/stringMinimiser/plots/animation.gif', format='GIF',
                        append_images=frames[1:] + frames[-2:0:-1],
                        save_all=True,duration=100,loop=0,subrectangles=True)



def evaluateAccuracy(variables:tuple, steps:np.ndarray, energyHistory_mean:np.ndarray, energyHistory_std:np.ndarray,
                     samplePositions:np.ndarray, compareIndex_pf:int, compareSampleIndex:int):
    """
    Evaluate the accuracy of the method.

    Inputs:
        - variables:tuple: Variables that were analysed.
        - steps:np.ndarray: MC steps.
        - energyHistory_mean:np.ndarray: History of the average energy.
        - energyHistory_std:np.ndarray: Standard deviations of the energy.
        - samplePositions:np.ndarray: Samples of positions of the samples (pointIndex,energyIndex,sample,dimension).
        - compareIndex_pf:int: pf sample to compare to.
        - compareSampleIndex:int: Node to use as comparison.
    """

    indicies = np.linspace(steps[0],steps[-1],samplePositions.shape[1], dtype=int)

    actualPositions = np.load(pathtohere / 'contents/stringMinimiser/data/testDataPoints.npy')

    differences_mean = np.ones(samplePositions.shape[1])
    differences_std = np.empty_like(differences_mean)

    for i in range(samplePositions.shape[1]):
        ds = np.abs(samplePositions[:,i,compareSampleIndex] - actualPositions)
        differences_mean[i] = np.mean(ds)
        differences_std[i] = np.std(ds)


    if len(indicies)!=len(differences_mean):
        print('len(indicies) != len(differences)')
        print(len(indicies),'!=',len(differences))
        raise Exception('Sizes of your arrays are not the same.')


    fig = plt.figure(figsize=(8,8))
    ax = fig.add_subplot()


    ax.errorbar(energyHistory_mean[indicies], differences_mean, marker='.',ls=None)
                    #xerr=energyHistory_std[log_indicies], yerr=differences_std)

    # Create appropiate labels.
    ax.set_xlabel('average energy')
    ax.set_ylabel('average difference')

    plt.savefig(pathtohere / 'contents/stringMinimiser/plots/accuracy.png', bbox_inches='tight')





def main():
    print('\nAnalysing data generated from string minimisation.')
    variables = ('log(ne)','log(t_e)','log(t_r)','log(h1)','log(xi)','log(he1)')
    
    # Load calculated sample positions.
    samplePositionsDims = np.loadtxt(
                    pathtohere / 'contents/stringMinimiser/data/samplePositionsDims.txt').T.astype(np.int64)
    samplePositions = np.loadtxt(pathtohere / 'contents/stringMinimiser/data/samplePositions.txt', delimiter=',',
                                    dtype=str)[:-1].reshape(samplePositionsDims).astype(np.float64)

    # Load known potentials.
    potentialPositionsDims = np.loadtxt(pathtohere / 'contents/stringMinimiser/data/potentialPositionsDims.txt').T.astype(np.int64)
    potentialPositions = np.loadtxt(pathtohere / 'contents/stringMinimiser/data/potentialPositions.txt',
                                    dtype=float).reshape(potentialPositionsDims)
    potentialPositions_pf = np.loadtxt(pathtohere / 'data/features.txt', delimiter=' ',
                                    dtype=float, skiprows=1,usecols=(1,2,3,4))
    # Load energy history.
    (steps,
     energyHistory_mean,
     energyHistory_std) = np.loadtxt(pathtohere / 'contents/stringMinimiser/data/energyHistory.txt', delimiter=',',
                                    dtype=float).T

    # Load actual positions.
    actualPositions = np.load(pathtohere / 'contents/stringMinimiser/data/testDataPoints.npy')
    


    plotEnergyHistory(steps, energyHistory_mean, energyHistory_std)


    pointIndex = -1
    compareIndex_pf = 15
    compareSampleIndex = 5

    print(samplePositions.shape)
    displayPairwise(variables, samplePositions[pointIndex,-1],
                    potentialPositions[:,pointIndex], actualPositions[pointIndex], compareSampleIndex)
    
    createFrames(samplePositions[:,-1].T, variables)
    

    print('\nGenerating PowerPoint')
    createPowerPoint(samplePositions.shape[2])

  

    print('\nComparing to actual run.')
    evaluateAccuracy(variables=tuple(map(lambda x: x[4:-1],variables)),steps=steps,
                        energyHistory_mean=energyHistory_mean, energyHistory_std=energyHistory_std,
                        samplePositions=samplePositions, compareIndex_pf=compareIndex_pf,
                        compareSampleIndex=compareSampleIndex)

    print('\nGenerating GIF')
    createGIF(samplePositions.shape[2])

    return 0


    

